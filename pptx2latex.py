import collections
import collections.abc

from pptx import Presentation
import pandas as pd
import aspose.slides as slides
from PIL import Image
import numpy as np

c = collections
c.abc = collections.abc

presentation = Presentation(
    r'C:\Users\faruk\OneDrive\Masaüstü\stock_watson_econ_4e_ppt_Ch01_02_03.pptx')


slide_template = r"""
\documentclass{beamer}
% Add necessary packages and customizations here

\usepackage{graphicx}
\setbeamertemplate{caption}[numbered]


\usetheme{CambridgeUS}

\title{Introduction to Econometrics} 
% \author{latex-beamer.com}
\date{\today}

\begin{document}

% Title slide
\begin{frame}
    \titlepage
\end{frame}

% Add more slides here

\end{document}
"""


def generate_slide_code(text_content, image_details, hasImage, table_content, is_indented):

    # Text blocks
    text_blocks = []
    if hasImage:
        text_block = r"""
    \begin{columns}
    \begin{column}{0.5\textwidth}
        \begin{itemize}
    """
    else:
        text_block = r"""
        \begin{itemize}
    """
    i = 0
    for text in text_content:

        if i != 0:
            if is_indented[i] == 1:
                text_block += r"""
                \begin{itemize}
                    \item %s
                \end{itemize}""" % text
            else:
                text_block += r"""
                \item %s
                """ % text
        i = i+1

    if hasImage:
        text_block += r"""
        \end{itemize}
        \end{column}"""
    else:
        text_block += r"""
        \end{itemize}"""
    text_blocks.append(text_block)

    # Image blocks
    image_blocks = []
    for image in image_details:  # %.2f
        if hasImage:
            image_block = r"""
        \begin{column}{0.5\textwidth}
            \begin{figure}
            \centering
                \includegraphics[width=0.7\textwidth]{%s}
                \caption{Image Caption}
            \end{figure}
        \end{column}
    \end{columns}
        """ % (image['filename'])  # image['width']

        else:
            image_block = r"""
        \begin{figure}
            \includegraphics[width=0.7\textwidth]{%s}
            \caption{Image Caption}
        \end{figure}
        """ % (image['filename'])  # image['width']
        image_blocks.append(image_block)

    table_code = ""
    if table_content:
        df = pd.DataFrame(table_content)
        df = df.fillna('')
        table_code_first = df.to_latex(
            index=False, header=False, bold_rows=True, caption="Table descripition")  # , header=False
        # first_line = f"\\begin{{tabular}}{{%s}}" % (len(table_content) * '|l')
        # print(table_code_first)
        table_code = r"""
        %s
        """ % (table_code_first)

    # Combine all elements into the slide code
    slide_code = r"""
    \begin{frame}
        \frametitle{%s}

        %s

        %s
        
        %s
    \end{frame}
    """ % (text_content[0], '\n'.join(text_blocks), '\n'.join(image_blocks), table_code)

    return slide_code


# Generate the Beamer code for each slide
beamer_code = []
for i, slide in enumerate(presentation.slides):
    text_content = []
    is_indented = []
    # indented_bullets = []
    image_details = []
    table_content = []
    hasImage = False

    # Retrieve the text content of the slide
    m = 0
    shape_number = []
    for shape in slide.shapes:
        # print('%d, %s' % (i, shape.shape_type))

        if shape.has_text_frame:
            m = m+1

            texts = []
            text_frame = shape.text_frame

            shape_number.append(m)
            for paragraph in text_frame.paragraphs:
                text_content.append(paragraph.text)
                if paragraph.level >= 1:  # paragraph.level >= 1
                    is_indented.append(1)
                else:
                    is_indented.append(2)

            # text = '\n'.join(p.text for p in text_frame.paragraphs)

            # text_content.append(text)

        elif shape.has_table:
            table = shape.table

            # Iterate over the rows and cells to retrieve table content

            for row in table.rows:
                row_content = []
                for cell in row.cells:
                    # print([p.text for p in cell.text_frame.paragraphs])
                    cell_text = [p.text for p in cell.text_frame.paragraphs]
                    if len(cell_text) > 1:
                        cell_text = cell_text[-1]
                    row_content.append(cell_text[0])
                # print("end of row")
                table_content.append(row_content)

    # Retrieve the image details of the slide
    for shape in slide.shapes:
        # print(i)
        # print(shape.shape_type)
        if shape.shape_type == 13:

            image = shape.image
            image_bytes = image.blob
            image_filename = fr'images/image_{i}_{shape.shape_id}.png'
            with open(image_filename, 'wb') as f:
                f.write(image_bytes)

            image_details.append({
                'filename': image_filename,
                'width': shape.width / 914400,  # Converting from EMU to inches
                'height': shape.height / 914400  # Converting from EMU to inches
            })

            hasImage = True

    # if i == 1 or i == 0:
    #     print(text_content)
    slide_code = generate_slide_code(
        text_content, image_details, hasImage, table_content, is_indented)
    beamer_code.append(slide_code)


# Join the slide codes and write to a .tex file
output_code = '\n'.join(beamer_code)
latex_code = slide_template.replace('% Add more slides here', output_code)
with open(r'C:\Users\faruk\OneDrive\Masaüstü\HIWI\latex\output.tex', 'w', encoding='utf-8') as f:
    f.write(latex_code)
